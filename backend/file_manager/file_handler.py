import os
import logging
from pathlib import Path
from typing import Tuple, Optional
import mimetypes
class FileHandler:
    logger = logging.getLogger(__name__)
    """Handles extraction of text from different file types."""

    SUPPORTED_TYPES = {
        # Text formats
        '.txt': 'text/plain',
        '.md': 'text/markdown',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.csv': 'text/csv',

        # Code files
        '.py': 'text/x-python',
        '.js': 'text/javascript',
        '.jsx': 'text/javascript',
        '.ts': 'text/typescript',
        '.tsx': 'text/typescript',
        '.java': 'text/x-java',
        '.cpp': 'text/x-c++src',
        '.c': 'text/x-csrc',
        '.h': 'text/x-chdr',
        '.cs': 'text/x-csharp',
        '.php': 'text/x-php',
        '.rb': 'text/x-ruby',
        '.go': 'text/x-go',
        '.rs': 'text/x-rust',
        '.swift': 'text/x-swift',
        '.kt': 'text/x-kotlin',
        '.scala': 'text/x-scala',
        '.sh': 'text/x-sh',
        '.bash': 'text/x-sh',
        '.sql': 'text/x-sql',
        '.html': 'text/html',
        '.css': 'text/css',
        '.scss': 'text/x-scss',
        '.yaml': 'text/yaml',
        '.yml': 'text/yaml',
        '.toml': 'text/x-toml',

        # Document formats
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',

        # Audio formats
        '.mp3': 'audio/mpeg',
        '.wav': 'audio/wav',
        '.m4a': 'audio/mp4',
        '.flac': 'audio/flac',
        '.ogg': 'audio/ogg',
        '.aac': 'audio/aac',

        # Video formats
        '.mp4': 'video/mp4',
        '.avi': 'video/x-msvideo',
        '.mov': 'video/quicktime',
        '.mkv': 'video/x-matroska',
        '.webm': 'video/webm',
        '.flv': 'video/x-flv',
    }
    
    @staticmethod
    def get_file_type(file_path: str) -> Optional[str]:
        """Get file type from path."""
        ext = Path(file_path).suffix.lower()
        return FileHandler.SUPPORTED_TYPES.get(ext)
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Optional[str]]:
        """
        Extract text from file.

        Args:
            file_path: Path to the file

        Returns:
            Tuple of (text_content, error_message)
            If successful, error_message is None
        """
        if not os.path.exists(file_path):
            return "", f"File not found: {file_path}"

        ext = Path(file_path).suffix.lower()

        # Text-based files (includes code, JSON, CSV, XML, YAML, etc.)
        if ext in ['.txt', '.md', '.json', '.xml', '.csv', '.py', '.js', '.jsx', '.ts',
                   '.tsx', '.java', '.cpp', '.c', '.h', '.cs', '.php', '.rb', '.go', '.rs',
                   '.swift', '.kt', '.scala', '.sh', '.bash', '.sql', '.html', '.css',
                   '.scss', '.yaml', '.yml', '.toml']:
            return FileHandler._extract_text_file(file_path)
        elif ext == '.pdf':
            return FileHandler._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return FileHandler._extract_docx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return FileHandler._extract_excel(file_path)
        elif ext in ['.pptx', '.ppt']:
            return FileHandler._extract_pptx(file_path)
        elif ext in ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac']:
            return FileHandler._extract_audio(file_path)
        elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv']:
            return FileHandler._extract_video(file_path)
        else:
            return "", f"Unsupported file type: {ext}"
    
    @staticmethod
    def _extract_text_file(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from TXT or MD files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content, None
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()
                return content, None
            except Exception as e:
                logger.error(f"Error reading text file {file_path}: {e}")
                return "", f"Error reading file: {str(e)}"
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return "", f"Error reading file: {str(e)}"
    
    @staticmethod
    def _extract_pdf(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from PDF file with fallback mechanisms for encoding issues."""
        # Suppress pdfplumber and pdfminer verbose logging during extraction
        pdfplumber_logger = logging.getLogger('pdfplumber')
        pdfminer_logger = logging.getLogger('pdfminer')
        old_pdfplumber_level = pdfplumber_logger.level
        old_pdfminer_level = pdfminer_logger.level
        
        try:
            pdfplumber_logger.setLevel(logging.WARNING)
            pdfminer_logger.setLevel(logging.WARNING)
            
            import pdfplumber
            
            text_parts = []
            
            # Try primary extraction method with proper encoding handling
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        try:
                            # Extract text with layout preservation
                            text = page.extract_text(layout=False)
                            
                            if text:
                                # Clean up encoding artifacts
                                text = FileHandler._clean_text(text)
                                if text.strip():  # Only add if there's content after cleaning
                                    text_parts.append(f"[Page {page_num}]\n{text}")
                            
                            # Fallback: Try extracting with layout if basic extraction gave no results
                            if not text or not text.strip():
                                text_layout = page.extract_text(layout=True)
                                if text_layout:
                                    text_layout = FileHandler._clean_text(text_layout)
                                    if text_layout.strip():
                                        text_parts.append(f"[Page {page_num}]\n{text_layout}")
                        
                        except Exception as page_error:
                            logger.warning(f"Error extracting page {page_num}: {page_error}")
                            # Continue to next page on error
                            continue
                
                if not text_parts:
                    return "", "No text content found in PDF"
                
                final_text = "\n\n".join(text_parts)
                
                # Verify the text is readable (not mostly garbage characters)
                if not FileHandler._is_valid_text(final_text):
                    logger.warning(f"PDF extraction resulted in mostly unreadable text, falling back to alternative method")
                    return FileHandler._extract_pdf_fallback(file_path)
                
                return final_text, None
            
            except Exception as e:
                logger.error(f"Error extracting PDF content with primary method: {e}")
                # Try fallback method
                return FileHandler._extract_pdf_fallback(file_path)
        
        except ImportError:
            logger.error("pdfplumber not installed")
            return "", "PDF support not available. Install pdfplumber."
        
        finally:
            # Restore original logging levels
            pdfplumber_logger.setLevel(old_pdfplumber_level)
            pdfminer_logger.setLevel(old_pdfminer_level)
    
    @staticmethod
    def _extract_pdf_fallback(file_path: str) -> Tuple[str, Optional[str]]:
        """Fallback PDF extraction using PyPDF2 if pdfplumber fails."""
        try:
            # Suppress PyPDF2 logging during import
            pdf2_logger = logging.getLogger('PyPDF2')
            old_level = pdf2_logger.level
            pdf2_logger.setLevel(logging.WARNING)
            
            try:
                from PyPDF2 import PdfReader
            finally:
                pdf2_logger.setLevel(old_level)
            
            text_parts = []
            
            try:
                pdf_reader = PdfReader(file_path)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        text = page.extract_text()
                        if text:
                            text = FileHandler._clean_text(text)
                            if text.strip():
                                text_parts.append(f"[Page {page_num}]\n{text}")
                    except Exception as page_error:
                        logger.warning(f"Error extracting page {page_num} with fallback: {page_error}")
                        continue
                
                if not text_parts:
                    return "", "No text content found in PDF (fallback method)"
                
                return "\n\n".join(text_parts), None
            
            except Exception as e:
                logger.error(f"Error with PyPDF2 fallback: {e}")
                return "", f"PDF extraction failed: {str(e)}"
        
        except ImportError:
            logger.error("PyPDF2 not installed for fallback PDF extraction")
            return "", "PDF extraction failed and fallback library not available"
    
    @staticmethod
    def _clean_text(text: str) -> str:
        """
        Clean extracted text by removing encoding artifacts and control characters.
        
        Args:
            text: Raw extracted text
            
        Returns:
            Cleaned text
        """
        if not text:
            return text
        
        # Remove null bytes and other control characters
        text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\r\t')
        
        # Fix common encoding artifacts
        # Replace mojibake patterns commonly found in PDF extraction
        replacements = [
            ('\x80', ''),  # NULL control character
            ('\x93', '"'),  # Left double quotation mark encoding issue
            ('\x94', '"'),  # Right double quotation mark encoding issue
            ('\x97', '—'),  # Em dash encoding issue
            ('\x96', '–'),  # En dash encoding issue
            ('\x92', "'"),  # Single quote encoding issue
        ]
        
        for old, new in replacements:
            text = text.replace(old, new)
        
        # Remove excessive whitespace
        lines = [line.rstrip() for line in text.split('\n')]
        text = '\n'.join(line for line in lines if line.strip())
        
        return text
    
    @staticmethod
    def _is_valid_text(text: str) -> bool:
        """
        Check if extracted text is valid and readable.
        
        Returns True if text has normal word patterns.
        Returns False if text is mostly control/encoding artifacts.
        
        Args:
            text: Text to validate
            
        Returns:
            Whether the text appears to be valid
        """
        if not text:
            return False
        
        # Clean the text first to remove obvious artifacts
        cleaned = FileHandler._clean_text(text)
        
        # Count words (sequences of alphanumeric + Unicode letter characters)
        import re
        words = re.findall(r'\b[a-zA-Z0-9\u0100-\uFFFF]+\b', cleaned)
        
        # Check average word length (real text typically has words 3+ chars)
        if words:
            avg_word_length = sum(len(w) for w in words) / len(words)
            # Text is valid if we have reasonable words and average word length is normal
            # (3-20 chars is typical for most languages)
            has_words = len(words) > 5  # At least a few words
            word_length_ok = 2 < avg_word_length < 25
            return has_words and word_length_ok
        
        # Fallback: check if we have mostly printable characters after cleaning
        printable_count = sum(1 for char in cleaned if ord(char) >= 32 or char in '\n\r\t')
        total_count = len(cleaned) if cleaned else len(text)
        
        if total_count == 0:
            return False
        
        # Text is valid if at least 60% is printable after cleaning
        readability = printable_count / total_count
        return readability >= 0.60

    @staticmethod
    def _extract_docx(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from DOCX or DOC files."""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables if any
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n\n".join(text_parts)
            
            if not text.strip():
                return "", "No text content found in DOCX"
            
            return text, None
        
        except ImportError:
            logger.error("python-docx not installed")
            return "", "DOCX support not available. Install python-docx."
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return "", f"Error processing DOCX: {str(e)}"

    @staticmethod
    def _extract_excel(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from XLSX or XLS files."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            if not text.strip():
                return "", "No text content found in Excel file"
            
            return text, None
        
        except ImportError:
            logger.error("openpyxl not installed")
            return "", "Excel support not available. Install openpyxl."
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {e}")
            return "", f"Error processing Excel: {str(e)}"

    @staticmethod
    def _extract_pptx(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from PPTX or PPT files."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[Slide {slide_num}]")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            text = "\n\n".join(text_parts)

            if not text.strip():
                return "", "No text content found in PowerPoint file"

            return text, None

        except ImportError:
            logger.error("python-pptx not installed")
            return "", "PowerPoint support not available. Install python-pptx."
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}")
            return "", f"Error processing PowerPoint: {str(e)}"

    @staticmethod
    def _extract_audio(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from audio files using speech recognition."""
        try:
            import speech_recognition as sr
            from pydub import AudioSegment
            import tempfile

            # Convert audio to WAV format if needed
            ext = Path(file_path).suffix.lower()

            if ext != '.wav':
                try:
                    logger.info(f"Converting {ext} audio to WAV format...")
                    audio = AudioSegment.from_file(file_path)

                    # Create temporary WAV file
                    with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_wav:
                        temp_wav_path = temp_wav.name
                        audio.export(temp_wav_path, format='wav')
                        working_file = temp_wav_path
                except Exception as e:
                    logger.error(f"Error converting audio format: {e}")
                    return "", f"Error converting audio format: {str(e)}"
            else:
                working_file = file_path

            try:
                # Use speech recognition to transcribe
                recognizer = sr.Recognizer()

                with sr.AudioFile(working_file) as source:
                    audio_data = recognizer.record(source)

                    try:
                        # Try Google Speech Recognition (free)
                        text = recognizer.recognize_google(audio_data)
                        return text, None
                    except sr.UnknownValueError:
                        return "", "Could not understand audio"
                    except sr.RequestError as e:
                        return "", f"Speech recognition service error: {str(e)}"
            finally:
                # Clean up temporary file if created
                if ext != '.wav' and os.path.exists(working_file):
                    try:
                        os.unlink(working_file)
                    except OSError:
                        pass  # Ignore cleanup errors

        except ImportError:
            logger.error("Audio processing libraries not installed")
            return "", "Audio support not available. Install speech_recognition and pydub packages."
        except Exception as e:
            logger.error(f"Error processing audio {file_path}: {e}")
            return "", f"Error processing audio: {str(e)}"

    @staticmethod
    def _extract_video(file_path: str) -> Tuple[str, Optional[str]]:
        """Extract text from video files by extracting audio and transcribing."""
        try:
            from moviepy.editor import VideoFileClip
            import tempfile

            logger.info(f"Extracting audio from video: {file_path}")

            # Extract audio from video
            try:
                video = VideoFileClip(file_path)

                # Create temporary audio file
                with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_audio:
                    temp_audio_path = temp_audio.name

                # Extract audio
                video.audio.write_audiofile(temp_audio_path, logger=None)
                video.close()

                # Use audio extraction method
                text, error = FileHandler._extract_audio(temp_audio_path)

                # Clean up
                try:
                    os.unlink(temp_audio_path)
                except OSError:
                    pass  # Ignore cleanup errors

                if error:
                    return "", f"Video processed but audio transcription failed: {error}"

                return text, None

            except Exception as e:
                logger.error(f"Error extracting audio from video: {e}")
                return "", f"Error extracting audio from video: {str(e)}"

        except ImportError:
            logger.error("Video processing libraries not installed")
            return "", "Video support not available. Install moviepy package."
        except Exception as e:
            logger.error(f"Error processing video {file_path}: {e}")
            return "", f"Error processing video: {str(e)}"


def extract_file_text(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Extract text from a file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Tuple of (text_content, error_message)
    """
    return FileHandler.extract_text(file_path)
